import os
import subprocess
import tkinter as tk
from tkinter import filedialog, messagebox
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("Lütfen 'pip install python-pptx' komutunu çalıştırın.")
    exit()

# ---------------- AYARLAR ----------------
# LibreOffice yolu
LIBREOFFICE_PATH = r"C:\Program Files\LibreOffice\program\soffice.exe"
# -----------------------------------------

def remove_media_from_pptx(input_path, output_path):
    """
    Sunum dosyasını açar, içindeki medya (video/ses) öğelerini siler
    ve yeni bir dosyaya kaydeder.
    """
    prs = Presentation(input_path)
    media_count = 0

    for slide in prs.slides:
        # Şekiller listesini kopyalayarak döngüye alıyoruz (silerken indeks kaymasın diye)
        for shape in list(slide.shapes):
            # Eğer şekil tipi MEDYA ise (Video veya Ses)
            if shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                # Şekli slayttan kaldır (XML ağacından siler)
                sp = shape._element
                sp.getparent().remove(sp)
                media_count += 1
    
    prs.save(output_path)
    return media_count

def convert_folder(folder_path, compress_mode):
    if not os.path.exists(LIBREOFFICE_PATH):
        messagebox.showerror("Hata", "LibreOffice bulunamadı! Kodun içindeki yolu kontrol edin.")
        return

    files = [f for f in os.listdir(folder_path) if f.endswith(".pptx")]
    
    if not files:
        messagebox.showwarning("Uyarı", "Klasörde .pptx dosyası bulunamadı. (.ppt dosyaları video silme özelliğini desteklemez)")
        return

    print(f"\n📂 Çalışılan Klasör: {folder_path}")
    print(f"⚙️ Sıkıştırma Modu (Video Silme): {'AÇIK' if compress_mode else 'KAPALI'}")
    print(f"📄 Toplam Dosya: {len(files)}\n")

    basarili = 0

    for i, filename in enumerate(files, 1):
        original_path = os.path.join(folder_path, filename)
        conversion_source = original_path # Varsayılan olarak orijinal dosyayı çevir
        temp_file_created = False
        
        print(f"[{i}/{len(files)}] İşleniyor: {filename}...")

        try:
            # Eğer sıkıştırma istenmişse
            if compress_mode:
                temp_filename = f"TEMP_CLEAN_{filename}"
                temp_path = os.path.join(folder_path, temp_filename)
                
                # Videoları silip geçici dosya oluştur
                deleted_count = remove_media_from_pptx(original_path, temp_path)
                
                if deleted_count > 0:
                    print(f"   ✂️  {deleted_count} adet video/medya silindi.")
                
                conversion_source = temp_path
                temp_file_created = True

            # LibreOffice ile PDF'e Çevir
            cmd = [
                LIBREOFFICE_PATH,
                "--headless",
                "--convert-to", "pdf",
                "--outdir", folder_path,
                conversion_source
            ]
            
            # Windows'ta konsol penceresini gizlemek için
            startupinfo = subprocess.STARTUPINFO()
            startupinfo.dwFlags |= subprocess.STARTF_USESHOWWINDOW

            process = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, startupinfo=startupinfo)

            if process.returncode == 0:
                print("   ✅ PDF oluşturuldu.")
                basarili += 1
            else:
                print("   ❌ LibreOffice Hatası.")

            # Temizlik: Geçici dosyayı sil (Eğer oluşturulduysa)
            if temp_file_created and os.path.exists(conversion_source):
                os.remove(conversion_source)
                # Orijinal dosyanın adıyla PDF oluştuğu için, TEMP_CLEAN isminde PDF kalırsa ismini düzeltelim
                temp_pdf = os.path.join(folder_path, f"TEMP_CLEAN_{os.path.splitext(filename)[0]}.pdf")
                target_pdf = os.path.join(folder_path, f"{os.path.splitext(filename)[0]}.pdf")
                
                if os.path.exists(temp_pdf):
                    # Eğer hedef PDF zaten varsa sil (üzerine yazmak için)
                    if os.path.exists(target_pdf):
                        os.remove(target_pdf)
                    os.rename(temp_pdf, target_pdf)

        except Exception as e:
            print(f"   ❌ Hata oluştu: {e}")
            # Hata durumunda da geçici dosyayı temizlemeyi dene
            if temp_file_created and os.path.exists(conversion_source):
                os.remove(conversion_source)

    messagebox.showinfo("Tamamlandı", f"İşlem Bitti.\nBaşarılı: {basarili}/{len(files)}")

def main():
    root = tk.Tk()
    root.withdraw()

    # 1. Klasör Seç
    folder_selected = filedialog.askdirectory(title="PPTX Dosyalarının Olduğu Klasörü Seç")
    if not folder_selected:
        return

    # 2. Sıkıştırma/Video Silme Sorusu
    cevap = messagebox.askyesno("Sıkıştırma Seçeneği", 
                                "Sunumların içindeki videolar silinsin mi?\n\n"
                                "EVET: Videolar silinir, dosya boyutu küçülür, PDF yapılır.\n"
                                "HAYIR: Videolar kalır (resim olarak görünür), boyut büyük olabilir.")

    convert_folder(folder_selected, compress_mode=cevap)

if __name__ == "__main__":
    main()